import pandas as pd
from pptx import Presentation
from pptx.util import Cm
from datetime import datetime
from PIL import Image, ImageDraw
import os

TEMPLATE_PATH = "template.pptx"
OUTPUT_FOLDER = "output/"
IMAGE_FOLDER = "images/"
TEMP_FOLDER = "images/temp/"

BIRTHDAY_SLIDE_INDEX = 5

os.makedirs(OUTPUT_FOLDER, exist_ok=True)
os.makedirs(TEMP_FOLDER, exist_ok=True)

df = pd.read_csv("employees.csv")

df['dob'] = pd.to_datetime(df['dob'], format='%d-%m-%Y', errors='coerce')
df['doj'] = pd.to_datetime(df['doj'], format='%d-%m-%Y', errors='coerce')

today = datetime.today()

def make_circle_image(input_path, output_path):
    img = Image.open(input_path).convert("RGBA")
    size = min(img.size)
    left = (img.width - size) // 2
    top = (img.height - size) // 2
    img = img.crop((left, top, left + size, top + size))

    mask = Image.new("L", (size, size), 0)
    draw = ImageDraw.Draw(mask)
    draw.ellipse((0, 0, size, size), fill=255)

    img.putalpha(mask)
    img.save(output_path)

def prepare_employee_image(emp_code):
    input_path = f"{IMAGE_FOLDER}{emp_code}.jpg"
    output_path = f"{TEMP_FOLDER}{emp_code}.png"

    if os.path.exists(input_path):
        make_circle_image(input_path, output_path)
        return output_path
    return None

def replace_name(slide, name):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if "{NAME}" in run.text:
                        run.text = run.text.replace("{NAME}", name)

def is_birthday(row):
    return row['dob'].day == today.day and row['dob'].month == today.month

def is_anniversary(row):
    return row['doj'].day == today.day and row['doj'].month == today.month

def get_years(row):
    return today.year - row['doj'].year

def remove_other_slides(prs, keep_index):
    slides_to_remove = []
    for i in range(len(prs.slides)):
        if i != keep_index:
            slides_to_remove.append(i)

    for i in reversed(slides_to_remove):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

def add_birthday_image(slide, img_path):
    if img_path:
        slide.shapes.add_picture(
            img_path,
            Cm(3.28),
            Cm(6.92),
            width=Cm(8.4),
            height=Cm(8.4)
        )

def add_year_image_junior(slide, years):
    if years in [5, 10]:
        return

    digits = list(str(years))

    if len(digits) == 1:
        path = f"images/years/Picture{digits[0]}.png"
        if os.path.exists(path):
            slide.shapes.add_picture(path, Cm(6.39), Cm(2.75), width=Cm(1.24), height=Cm(2.15))

    elif len(digits) == 2:
        positions = [(5.9, 2.6), (6.88, 2.58)]
        for i, digit in enumerate(digits):
            path = f"images/years/Picture{digit}.png"
            if os.path.exists(path):
                slide.shapes.add_picture(path, Cm(positions[i][0]), Cm(positions[i][1]), width=Cm(1.24), height=Cm(2.15))

def add_year_image_senior(slide, years):
    if years in [5, 10]:
        return

    digits = list(str(years))

    if len(digits) == 1:
        path = f"images/years/Picture{digits[0]}.png"
        if os.path.exists(path):
            slide.shapes.add_picture(path, Cm(6.24), Cm(2.58), width=Cm(1.24), height=Cm(2.15))

    elif len(digits) == 2:
        positions = [(5.6, 2.53), (6.58, 2.53)]
        for i, digit in enumerate(digits):
            path = f"images/years/Picture{digit}.png"
            if os.path.exists(path):
                slide.shapes.add_picture(path, Cm(positions[i][0]), Cm(positions[i][1]), width=Cm(1.24), height=Cm(2.15))

def add_senior_image(slide, img_path):
    if img_path:
        pic = slide.shapes.add_picture(
            img_path,
            Cm(9.27),
            Cm(3.87),
            width=Cm(11.87),
            height=Cm(11.76)
        )
        return pic

def send_to_back(shape):
    sp = shape._element
    parent = sp.getparent()
    parent.remove(sp)
    parent.insert(0, sp)

for _, row in df.iterrows():

    if is_birthday(row):
        prs = Presentation(TEMPLATE_PATH)
        slide = prs.slides[BIRTHDAY_SLIDE_INDEX]

        replace_name(slide, row['name'])
        img_path = prepare_employee_image(row['employee_code'])
        add_birthday_image(slide, img_path)

        remove_other_slides(prs, BIRTHDAY_SLIDE_INDEX)

        prs.save(f"{OUTPUT_FOLDER}birthday_{row['employee_code']}.pptx")

    if is_anniversary(row):

        years = get_years(row)
        prs = Presentation(TEMPLATE_PATH)

        if row['grade'] in ["EVP", "Partner"]:

            if years == 5:
                slide_index = 2
            elif years == 10:
                slide_index = 3
            else:
                slide_index = 1

            slide = prs.slides[slide_index]

            replace_name(slide, row['name'])

            img_path = prepare_employee_image(row['employee_code'])
            pic = add_senior_image(slide, img_path)
            send_to_back(pic)

            if years not in [5, 10]:
                add_year_image_senior(slide, years)

            remove_other_slides(prs, slide_index)

            prs.save(f"{OUTPUT_FOLDER}anniv_senior_{row['employee_code']}.pptx")

        elif row['grade'] in ["JM", "Analyst", "Manager", "Senior Manager"]:

            if years == 5:
                slide_index = 2
            elif years == 10:
                slide_index = 3
            else:
                slide_index = 0

            slide = prs.slides[slide_index]

            replace_name(slide, row['name'])

            img_path = f"{IMAGE_FOLDER}{row['employee_code']}.jpg"

            if years == 5:
                slide.shapes.add_picture(img_path, Cm(17.41), Cm(3.84), width=Cm(7.32), height=Cm(8.03))
            elif years == 10:
                slide.shapes.add_picture(img_path, Cm(16.13), Cm(4.71), width=Cm(7.37), height=Cm(7.37))
            else:
                slide.shapes.add_picture(img_path, Cm(16.33), Cm(4.67), width=Cm(8.5), height=Cm(9.45))

            if years not in [5, 10]:
                add_year_image_junior(slide, years)

            remove_other_slides(prs, slide_index)

            prs.save(f"{OUTPUT_FOLDER}anniv_{row['employee_code']}.pptx")